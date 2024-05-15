import streamlit as st
import glob
import os
import random
import re
import string
import openai
from icrawler.builtin import GoogleImageCrawler
from uuid import uuid4
from pptx import Presentation

def generate_random_string():
    return ''.join(random.choice(string.ascii_uppercase + string.ascii_lowercase + string.digits) for _ in range(16))

random_prefix = generate_random_string()

class PrefixNameDownloader(GoogleImageCrawler):
    
    def get_filename(self, task, default_ext):
        filename = super().get_filename(task, default_ext)
        return f'prefix_{random_prefix}_{filename}'

def generate_ppt(theme, topic, slide_length, api_key):
    root = Presentation(theme)
    openai.api_key = api_key

    message_input = f"""
    Create content for a slideshow presentation.
    The content's topic is {topic}. 
    The slideshow is {slide_length} slides long. 
    The content is written in the language of the content I give you above.
    
    
    You are allowed to use the following slide types:
    
    Slide types:
    Title Slide - (Title, Subtitle)
    Content Slide - (Title, Content)
    Image Slide - (Title, Content, Image)
    Thanks Slide - (Title)
    
    Put this tag before the Title Slide: [L_TS]
    Put this tag before the Content Slide: [L_CS]
    Put this tag before the Image Slide: [L_IS]
    Put this tag before the Thanks Slide: [L_THS]
    
    Put "[SLIDEBREAK]" after each slide 
    
    For example:
    [L_TS]
    [TITLE]Mental Health[/TITLE]
    
    [SLIDEBREAK]
    
    [L_CS] 
    [TITLE]Mental Health Definition[/TITLE]
    [CONTENT]
    1. Definition: A person’s condition with regard to their psychological and emotional well-being
    2. Can impact one's physical health
    3. Stigmatized too often.
    [/CONTENT]
    
    [SLIDEBREAK]
    
    Put this tag before the Title: [TITLE]
    Put this tag after the Title: [/TITLE]
    Put this tag before the Subitle: [SUBTITLE]
    Put this tag after the Subtitle: [/SUBTITLE]
    Put this tag before the Content: [CONTENT]
    Put this tag after the Content: [/CONTENT]
    Put this tag before the Image: [IMAGE]
    Put this tag after the Image: [/IMAGE]
    
    Elaborate on the Content, provide as much information as possible.
    You put a [/CONTENT] at the end of the Content.
    Do not reply as if you are talking about the slideshow itself. (ex. "Include pictures here about...")
    Do not include any special characters (?, !, ., :, ) in the Title.
    Do not include any additional information in your response and stick to the format."""

   
    completion = openai.chat.completions.create(
        model="gpt-4",
        messages=[
            {
                "role": "user",
                "content": message_input,
            },
        ],
    )
    parsed_content = completion.choices[0].message.content
  


    # """ Ref for slide types:
    # 0 -> title and subtitle
    # 1 -> title and content
    # 2 -> section header
    # 3 -> two content
    # 4 -> Comparison
    # 5 -> Title only
    # 6 -> Blank
    # 7 -> Content with caption
    # 8 -> Pic with caption
    # """

    def delete_all_slides():
        for i in range(len(root.slides) - 1, -1, -1):
            r_id = root.slides._sldIdLst[i].rId
            root.part.drop_rel(r_id)
            del root.slides._sldIdLst[i]

    def create_title_slide(title, subtitle):
        layout = root.slide_layouts[0]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    def create_section_header_slide(title):
        layout = root.slide_layouts[2]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = title

    def create_title_and_content_slide(title, content):
        layout = root.slide_layouts[1]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = content

    def create_title_and_content_and_image_slide(title, content, image_query):
        layout = root.slide_layouts[8]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[2].text = content
        refresh_bad_coding_practice()
        bing_crawler = GoogleImageCrawler(downloader_cls=PrefixNameDownloader, storage={'root_dir': os.getcwd()})
        bing_crawler.crawl(keyword=image_query, max_num=1)
        dir_path = os.path.dirname(os.path.realpath(__file__))
        file_name = glob.glob(f"prefix_{bad_coding_practice}*")
        print(file_name)
        img_path = os.path.join(dir_path, file_name[0])
        slide.shapes.add_picture(img_path, slide.placeholders[1].left, slide.placeholders[1].top,
                                 slide.placeholders[1].width, slide.placeholders[1].height)

    def find_text_in_between_tags(text, start_tag, end_tag):
        start_pos = text.find(start_tag)
        end_pos = text.find(end_tag)
        result = []
        while start_pos > -1 and end_pos > -1:
            text_between_tags = text[start_pos + len(start_tag):end_pos]
            result.append(text_between_tags)
            start_pos = text.find(start_tag, end_pos + len(end_tag))
            end_pos = text.find(end_tag, start_pos)
        res1 = "".join(result)
        res2 = re.sub(r"\[IMAGE\].*?\[/IMAGE\]", '', res1)
        if len(result) > 0:
            return res2
        else:
            return ""

    def search_for_slide_type(text):
        tags = ["[L_TS]", "[L_CS]", "[L_IS]", "[L_THS]"]
        found_text = next((s for s in tags if s in text), None)
        return found_text

    def parse_response(reply):
        list_of_slides = reply.split("[SLIDEBREAK]")
        for slide in list_of_slides:
            slide_type = search_for_slide_type(slide)
            if slide_type == "[L_TS]":
                create_title_slide(find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]"),
                                   find_text_in_between_tags(str(slide), "[SUBTITLE]", "[/SUBTITLE]"))
            elif slide_type == "[L_CS]":
                create_title_and_content_slide("".join(find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")),
                                               "".join(find_text_in_between_tags(str(slide), "[CONTENT]",
                                                                                 "[/CONTENT]")))
            elif slide_type == "[L_IS]":
                create_title_and_content_and_image_slide("".join(find_text_in_between_tags(str(slide), "[TITLE]",
                                                                                           "[/TITLE]")),
                                                         "".join(find_text_in_between_tags(str(slide), "[CONTENT]",
                                                                                           "[/CONTENT]")),
                                                         "".join(find_text_in_between_tags(str(slide), "[IMAGE]",
                                                                                           "[/IMAGE]")))
            elif slide_type == "[L_THS]":
                create_section_header_slide("".join(find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")))

    def find_title():
        return root.slides[0].shapes.title.text

    delete_all_slides()

    # print(response)

    # parse_response(response['choices'][0]['message']['content'])

    name_ = str(uuid4()).replace('-', '')

    root.save(f"./{name_}.pptx")

    print("done")

    dir_path = "./"
    prefix = "prefix_"

    for file_name in os.listdir(dir_path):
        if file_name.startswith(prefix):
            file_path = os.path.join(dir_path, file_name)
            if os.path.isfile(file_path):
                os.remove(file_path)

    return f"./{name_}.pptx"

    # parsed_content = response['choices'][0]['message']['content']
    
    # Assume more functions here for processing the content into slides
    # Functions like create_title_slide, create_content_slide, etc., as needed

    ppt_filename = str(uuid4()).replace('-', '') + '.pptx'
    root.save(ppt_filename)

    # Cleanup any downloaded images
    for file_name in os.listdir("."):
        if file_name.startswith(f'prefix_{random_prefix}_'):
            os.remove(file_name)

    return ppt_filename

st.title("ChatGPT PPT Framework Generator")

api_key = st.text_input("Enter your OpenAI API Key")
topic = st.text_input("Enter the topic or content for the PPT")
slide_length = st.slider("Select the number of slides for the PPT", min_value=1, max_value=20, value=6, step=1)
theme = st.file_uploader("Upload your PPT template", type=['pptx', 'ppt'])

if st.button("Generate PPT"):
    if theme is not None and topic and api_key:
        ppt_path = generate_ppt(theme, topic, slide_length, api_key)
        st.success("Generated PPT successfully!")
        with open(ppt_path, "rb") as file:
            st.download_button(
                label="Download PPT",
                data=file,
                file_name="Generated_Presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
    else:
        st.error("Please fill all inputs and upload a file.")

